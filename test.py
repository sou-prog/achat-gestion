import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import os
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from datetime import datetime
from st_aggrid import AgGrid, GridOptionsBuilder
import base64
from io import BytesIO
from sklearn.linear_model import LinearRegression
import numpy as np
import sqlite3
from pptx import Presentation
from pptx.util import Inches
from supabase import create_client, Client
from dotenv import load_dotenv
import uuid

# Load environment variables from .env file
load_dotenv()

# Check required environment variables
required_env_vars = ["SUPABASE_URL", "SUPABASE_KEY"]
missing_vars = [var for var in required_env_vars if not os.getenv(var)]
if missing_vars:
    st.error(f"Variables d'environnement manquantes : {', '.join(missing_vars)}. Vérifiez le fichier .env.")
    st.stop()

# Check optional SMTP variables
smtp_vars = ["SMTP_SERVER", "SMTP_PORT", "SMTP_USERNAME", "SMTP_PASSWORD", "NOTIFICATION_RECIPIENT"]
smtp_available = all(os.getenv(var) for var in smtp_vars)

# Initialize Supabase client
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_KEY")
supabase: Client = create_client(SUPABASE_URL, SUPABASE_KEY)

# Configure page
st.set_page_config(page_title="Indirect Purchases Dashboard", layout="wide", initial_sidebar_state="expanded")

# Define theme styles
def set_theme(theme):
    if theme == "Dark":
        st.markdown(
            """
            <style>
            .stApp {
                background-color: #1E1E1E;
                color: #FFFFFF;
                font-family: 'Arial', sans-serif;
            }
            .css-1aumxhk {
                background-color: #2E2E2E;
            }
            .stButton>button {
                background-color: #005BAC;
                color: white;
                border-radius: 5px;
            }
            .stButton>button:hover {
                background-color: #003D7A;
            }
            .stHeader {
                font-size: 24px;
                font-weight: bold;
                color: #FFFFFF !important;
                text-shadow: 1px 1px 2px #000000;
            }
            .stSubheader {
                font-size: 18px;
                font-weight: bold;
                color: #F0F0F0 !important;
                text-shadow: 1px 1px 1px #000000;
            }
            .stWarning {
                color: #FF4B4B;
            }
            .section {
                padding: 20px;
                background-color: #252525;
                border-radius: 10px;
                margin-bottom: 20px;
            }
            .stMetric {
                color: #FFFFFF !important;
            }
            .stMetric [data-testid="stMetricLabel"] {
                color: #FFFFFF !important;
                font-weight: bold;
            }
            .stMetric [data-testid="stMetricValue"] {
                color: #00FF7F !important;
                font-size: 20px;
            }
            .stTabs [data-testid="stTabs"] {
                background-color: #2E2E2E;
            }
            .stTabs [role="tab"] {
                color: #FFFFFF !important;
                font-weight: bold;
            }
            .stTabs [role="tab"][aria-selected="true"] {
                background-color: #005BAC !important;
                color: #FFFFFF !important;
            }
            .stTabs [role="tab"]:hover {
                background-color: #003D7A !important;
                color: #FFFFFF !important;
            }
            </style>
            """,
            unsafe_allow_html=True
        )
    else:
        st.markdown(
            """
            <style>
            .stApp {
                background-color: #F5F5F5;
                color: #333333;
                font-family: 'Arial', sans-serif;
            }
            .css-1aumxhk {
                background-color: #E0E0E0;
            }
            .stButton>button {
                background-color: #005BAC;
                color: white;
                border-radius: 5px;
            }
            .stButton>button:hover {
                background-color: #003D7A;
            }
            .stHeader {
                font-size: 24px;
                font-weight: bold;
                color: #0073E6 !important;
                text-shadow: 1px 1px 2px #AAAAAA;
            }
            .stSubheader {
                font-size: 18px;
                font-weight: bold;
                color: #4D4D4D !important;
                text-shadow: 1px 1px 1px #CCCCCC;
            }
            .stWarning {
                color: #CC0000;
            }
            .section {
                padding: 20px;
                background-color: #FFFFFF;
                border-radius: 10px;
                margin-bottom: 20px;
                box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
            }
            .stMetric {
                color: #333333 !important;
            }
            .stMetric [data-testid="stMetricLabel"] {
                color: #333333 !important;
                font-weight: bold;
            }
            .stMetric [data-testid="stMetricValue"] {
                color: #0073E6 !important;
                font-size: 20px;
            }
            .stTabs [data-testid="stTabs"] {
                background-color: #E0E0E0;
            }
            .stTabs [role="tab"] {
                color: #333333 !important;
                font-weight: bold;
            }
            .stTabs [role="tab"][aria-selected="true"] {
                background-color: #005BAC !important;
                color: #FFFFFF !important;
            }
            .stTabs [role="tab"]:hover {
                background-color: #003D7A !important;
                color: #FFFFFF !important;
            }
            </style>
            """,
            unsafe_allow_html=True
        )

# Translations
translations = {
    "fr": {
        "title": "Tableau de Bord Achats Indirects",
        "loading": "Chargement des données... 📊",
        "data_loaded": "Données chargées avec succès ! ✅",
        "summary": "Résumé Global 🌐",
        "filters_alerts": "Filtres et Alertes ⚙️",
        "po_filters": "Filtres Bons de Commande 📋",
        "pt_filters": "Filtres Conditions de Paiement 💰",
        "contract_filters": "Filtres Contrats 📜",
        "supplier": "Fournisseur",
        "department": "Département",
        "purchase_type": "Type d'achat",
        "status": "Statut",
        "division": "Division",
        "period": "Période",
        "amount_threshold": "Seuil d'alerte montant (EUR)",
        "delay_threshold": "Seuil d'alerte retard paiement (jours)",
        "alerts": "Alertes 🚨",
        "no_alerts": "Aucune alerte active.",
        "check_alerts": "Vérifier alertes",
        "po_tab": "Bons de Commande",
        "pt_tab": "Conditions de Paiement",
        "contract_tab": "Contrats",
        "po_header": "Suivi des Bons de Commande 📦",
        "pt_header": "Suivi des Conditions de Paiement 💸",
        "contract_header": "Suivi des Contrats 📜",
        "no_data": "Aucune donnée disponible pour les filtres sélectionnés.",
        "po_by_dept": "Commandes par département 📊",
        "amount_quantity": "Évolution montants et quantités 📈",
        "status_dist": "Répartition par statut",
        "type_dist": "Répartition par type d'achat",
        "new_terms": "Nouveaux termes de paiement",
        "old_terms": "Anciens termes de paiement",
        "kpis": "Indicateurs de Performance 📋",
        "turnover": "Chiffre d'affaires (EUR)",
        "improvement": "Amélioration (%)",
        "cash_flow": "Gain trésorerie (EUR)",
        "terms_by_division": "Délais par division",
        "forecast": "Prévisions des Achats 📈",
        "predict_by": "Prédire par",
        "compare_fournisseurs": "Comparaison des Fournisseurs 🌟",
        "heatmap": "Performances Fournisseur-Division 🌡️",
        "reorder": "Suggestions de Réapprovisionnement 🛒",
        "reorder_threshold": "Seuil de réapprovisionnement (quantité)",
        "comments": "Commentaires",
        "add_comment": "Ajouter commentaire",
        "comment_text": "Commentaire",
        "comment_user": "Utilisateur",
        "footer": "Développé pour KOSTAL | Mis à jour : ",
        "export_chart": "Exporter graphique en PNG",
        "export_ppt": "Exporter en PowerPoint",
        "send_contract_reminders": "Envoyer rappels contrats",
        "select_filter": "Veuillez sélectionner au moins une option pour chaque filtre.",
        "total_orders": "Total des commandes",
        "pending_orders": "Commandes en attente",
        "total_turnover": "Chiffre d'affaires total",
        "help": "Aide ℹ️",
        "help_text": "Utilisez les filtres latéraux pour affiner les données. Cliquez sur 'Vérifier alertes' pour envoyer des notifications. Exportez les graphiques via les boutons appropriés.",
        "theme": "Thème",
        "color_scheme": "Palette de couleurs",
        "login": "Connexion",
        "email": "Email",
        "password": "Mot de passe",
        "login_button": "Se connecter",
        "logout_button": "Se déconnecter",
        "login_success": "Connexion réussie !",
        "login_error": "Erreur de connexion : email ou mot de passe incorrect.",
        "please_login": "Veuillez vous connecter pour accéder au tableau de bord."
    },
    "en": {
        "title": "Indirect Purchases Dashboard",
        "loading": "Loading data... 📊",
        "data_loaded": "Data loaded successfully! ✅",
        "summary": "Global Summary 🌐",
        "filters_alerts": "Filters and Alerts ⚙️",
        "po_filters": "Purchase Orders Filters 📋",
        "pt_filters": "Payment Terms Filters 💰",
        "contract_filters": "Contracts Filters 📜",
        "supplier": "Supplier",
        "department": "Department",
        "purchase_type": "Purchase Type",
        "status": "Status",
        "division": "Division",
        "period": "Period",
        "amount_threshold": "Amount alert threshold (EUR)",
        "delay_threshold": "Payment delay alert threshold (days)",
        "alerts": "Alerts 🚨",
        "no_alerts": "No active alerts.",
        "check_alerts": "Check alerts",
        "po_tab": "Purchase Orders",
        "pt_tab": "Payment Terms",
        "contract_tab": "Contracts",
        "po_header": "Purchase Orders Tracking 📦",
        "pt_header": "Payment Terms Tracking 💸",
        "contract_header": "Contracts Tracking 📜",
        "no_data": "No data available for the selected filters.",
        "po_by_dept": "Orders by Department 📊",
        "amount_quantity": "Amount and Quantity Evolution 📈",
        "status_dist": "Distribution by Status",
        "type_dist": "Distribution by Purchase Type",
        "new_terms": "New Payment Terms",
        "old_terms": "Old Payment Terms",
        "kpis": "Key Performance Indicators 📋",
        "turnover": "Turnover (EUR)",
        "improvement": "Improvement (%)",
        "cash_flow": "Cash Flow Gain (EUR)",
        "terms_by_division": "Terms by Division",
        "forecast": "Purchase Forecasts 📈",
        "predict_by": "Predict by",
        "compare_fournisseurs": "Suppliers Comparison 🌟",
        "heatmap": "Supplier-Division Performance 🌡️",
        "reorder": "Reorder Suggestions 🛒",
        "reorder_threshold": "Reorder threshold (quantity)",
        "comments": "Comments",
        "add_comment": "Add comment",
        "comment_text": "Comment",
        "comment_user": "User",
        "footer": "Developed for KOSTAL | Updated: ",
        "export_chart": "Export chart as PNG",
        "export_ppt": "Export to PowerPoint",
        "send_contract_reminders": "Send contract reminders",
        "select_filter": "Please select at least one option for each filter.",
        "total_orders": "Total Orders",
        "pending_orders": "Pending Orders",
        "total_turnover": "Total Turnover",
        "help": "Help ℹ️",
        "help_text": "Use the sidebar filters to refine data. Click 'Check alerts' to send notifications. Export charts using the appropriate buttons.",
        "theme": "Theme",
        "color_scheme": "Color scheme",
        "login": "Login",
        "email": "Email",
        "password": "Password",
        "login_button": "Log In",
        "logout_button": "Log Out",
        "login_success": "Login successful!",
        "login_error": "Login error: incorrect email or password.",
        "please_login": "Please log in to access the dashboard."
    }
}

# Language selector
lang = st.sidebar.selectbox("Langue / Language", ["Français", "English"], key="lang_select")
lang_key = "fr" if lang == "Français" else "en"
t = translations[lang_key]

# Authentication management
if "logged_in" not in st.session_state:
    st.session_state.logged_in = False
    st.session_state.user_email = None
    st.session_state.login_message = None

# Login screen
with st.sidebar:
    st.subheader(t["login"])
    login_placeholder = st.empty()
    if not st.session_state.logged_in:
        email = st.text_input(t["email"], key="login_email")
        password = st.text_input(t["password"], type="password", key="login_password")
        if st.button(t["login_button"], key="login_btn"):
            if not email or not password:
                login_placeholder.error("Veuillez entrer un email et un mot de passe.")
            else:
                with st.spinner("Connexion en cours..."):
                    try:
                        response = supabase.auth.sign_in_with_password({"email": email, "password": password})
                        if response.user:
                            st.session_state.logged_in = True
                            st.session_state.user_email = email
                            login_placeholder.success(t["login_success"])
                        else:
                            login_placeholder.error(t["login_error"])
                    except Exception as e:
                        login_placeholder.error(f"Erreur de connexion : {str(e)}")
    else:
        login_placeholder.write(f"Connecté en tant que : {st.session_state.user_email}")
        if st.button(t["logout_button"], key="logout_btn"):
            try:
                supabase.auth.sign_out()
                st.session_state.logged_in = False
                st.session_state.user_email = None
                login_placeholder.success("Déconnexion réussie !")
            except Exception as e:
                login_placeholder.error(f"Erreur lors de la déconnexion : {str(e)}")

# Check if user is logged in
if not st.session_state.logged_in:
    st.warning(t["please_login"])
    st.stop()

# Theme selector
theme = st.sidebar.selectbox(t["theme"], ["Dark", "Light"], key="theme_select")
set_theme(theme)

# Color scheme selector
color_scheme = st.sidebar.selectbox(t["color_scheme"], ["Plotly", "Viridis", "Cividis", "Inferno"], key="color_scheme")
color_schemes = {
    "Plotly": px.colors.qualitative.Plotly,
    "Viridis": px.colors.sequential.Viridis,
    "Cividis": px.colors.sequential.Cividis,
    "Inferno": px.colors.sequential.Inferno
}

# Send email function
def send_email(subject, body, to_email):
    if not smtp_available:
        st.error("⚠️ Configuration SMTP manquante. Vérifiez le fichier .env.")
        return False
    smtp_server = os.getenv("SMTP_SERVER")
    smtp_port = os.getenv("SMTP_PORT")
    smtp_username = os.getenv("SMTP_USERNAME")
    smtp_password = os.getenv("SMTP_PASSWORD")

    msg = MIMEMultipart()
    msg['From'] = smtp_username
    msg['To'] = to_email
    msg['Subject'] = subject
    msg.attach(MIMEText(body, 'plain'))

    try:
        with smtplib.SMTP(smtp_server, int(smtp_port)) as server:
            server.starttls()
            server.login(smtp_username, smtp_password)
            server.sendmail(smtp_username, to_email, msg.as_string())
        return True
    except Exception as e:
        st.error(f"⚠️ Erreur lors de l'envoi de l'email : {str(e)}")
        return False

# Check alerts function
def check_alerts(df_contracts, df_po, df_pt):
    if not smtp_available:
        st.error("⚠️ Destinataire des notifications ou configuration SMTP non configuré dans le fichier .env.")
        return 0
    to_email = os.getenv("NOTIFICATION_RECIPIENT")
    notifications_sent = 0

    for index, row in df_contracts.iterrows():
        days_left = (row["DATE_EXPIRATION"] - pd.Timestamp.now()).days
        if days_left <= 60:
            subject = f"Alert: {t['status']} {row['CONTRAT']} Expiration"
            body = (
                f"{t['status']} {row['CONTRAT']} with {row['FOURNISSEUR']} "
                f"expires in {days_left} days.\n"
                f"Expiration date: {row['DATE_EXPIRATION'].strftime('%d/%m/%Y')}\n"
                f"Amount: {row['MONTANT_MAD']:,.2f} MAD"
            )
            if send_email(subject, body, to_email):
                notifications_sent += 1

    for index, row in df_po[df_po["STATUT"] == "En attente"].iterrows():
        subject = f"Alert: {t['status']} {row['PO_NUMBER']} Pending"
        body = (
            f"{t['status']} {row['PO_NUMBER']} with {row['FOURNISSEUR']} "
            f"is pending validation.\n"
            f"Amount: {row['MONTANT_EUR']:,.2f} EUR\n"
            f"Department: {row['DEPARTEMENT']}\n"
            f"Date: {row['DATE'].strftime('%d/%m/%Y')}"
        )
        if send_email(subject, body, to_email):
            notifications_sent += 1

    for index, row in df_pt[df_pt["DELAI_PAIEMENT"] > 0].iterrows():
        subject = f"Alert: {t['status']} {row['FOURNISSEUR']} Payment Delay"
        body = (
            f"{t['status']} {row['FOURNISSEUR']} has a payment delay of {row['DELAI_PAIEMENT']} days.\n"
            f"Turnover: {row['TURNOVER_EUR']:,.2f} EUR\n"
            f"Division: {row['DIVISION']}"
        )
        if send_email(subject, body, to_email):
            notifications_sent += 1

    return notifications_sent

# Load data from Supabase
@st.cache_data
def load_data(table_name, required_cols, _placeholder=None):
    try:
        response = supabase.table(table_name).select("*").execute()
        if not response.data:
            return pd.DataFrame(columns=required_cols), f"Aucune donnée trouvée dans la table {table_name}. Vérifiez si la table existe et contient des données."
        df = pd.DataFrame(response.data)
        actual_cols = df.columns.tolist()
        column_mapping = {
            'po_number': 'PO_NUMBER',
            'fournisseur': 'FOURNISSEUR',
            'departement': 'DEPARTEMENT',
            'montant_eur': 'MONTANT_EUR',
            'quantite': 'QUANTITE',
            'date': 'DATE',
            'type_achat': 'TYPE_ACHAT',
            'statut': 'STATUT',
            'old_days': 'OLD_DAYS',
            'new_days': 'NEW_DAYS',
            'turnover_eur': 'TURNOVER_EUR',
            'division': 'DIVISION',
            'condition_paiement': 'CONDITION_PAIEMENT',
            'delai_paiement': 'DELAI_PAIEMENT',
            'contrat': 'CONTRAT',
            'date_expiration': 'DATE_EXPIRATION',
            'montant_mad': 'MONTANT_MAD',
            'responsable_email': 'RESPONSABLE_EMAIL'
        }
        df.rename(columns={k: v for k, v in column_mapping.items() if k in df.columns}, inplace=True)
        missing_cols = [col for col in required_cols if col not in df.columns and col.lower() not in [c.lower() for c in df.columns]]
        if missing_cols:
            return pd.DataFrame(columns=required_cols), f"Colonnes manquantes dans {table_name}: {missing_cols}. Colonnes disponibles: {actual_cols}"
        for col in df.columns:
            for req_col in required_cols:
                if col.lower() == req_col.lower() and col != req_col:
                    df.rename(columns={col: req_col}, inplace=True)
        return df, None
    except Exception as e:
        return pd.DataFrame(columns=required_cols), f"⚠️ Erreur lors du chargement de {table_name}: {str(e)}"

# Export Plotly figure as PNG
def export_plotly_figure(fig, filename):
    if fig is not None:
        img_bytes = fig.to_image(format="png")
        b64 = base64.b64encode(img_bytes).decode()
        href = f'<a href="data:image/png;base64,{b64}" download="{filename}.png">{t["export_chart"]}</a>'
        st.markdown(href, unsafe_allow_html=True)
    else:
        st.warning("⚠️ Aucune figure disponible pour l'exportation.")

# Export to PowerPoint
def export_to_ppt(df_po_filtered, df_pt_filtered, df_contracts_filtered, figs):
    prs = Presentation()
    
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = t["po_header"]
    
    for fig_name, fig in figs.items():
        if fig:
            img_bytes = fig.to_image(format="png")
            img_stream = BytesIO(img_bytes)
            slide.shapes.add_picture(img_stream, Inches(1), Inches(1.5), width=Inches(8))
    
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = t["po_tab"]
    rows, cols = df_po_filtered.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(1), Inches(1.5), Inches(8), Inches(3)).table
    for j, col in enumerate(df_po_filtered.columns):
        table.cell(0, j).text = col
    for i in range(rows):
        for j in range(cols):
            table.cell(i + 1, j).text = str(df_po_filtered.iloc[i, j])
    
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = t["pt_tab"]
    rows, cols = df_pt_filtered.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(1), Inches(1.5), Inches(8), Inches(3)).table
    for j, col in enumerate(df_pt_filtered.columns):
        table.cell(0, j).text = col
    for i in range(rows):
        for j in range(cols):
            table.cell(i + 1, j).text = str(df_pt_filtered.iloc[i, j])
    
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = t["contract_tab"]
    rows, cols = df_contracts_filtered.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(1), Inches(1.5), Inches(8), Inches(3)).table
    for j, col in enumerate(df_contracts_filtered.columns):
        table.cell(0, j).text = col
    for i in range(rows):
        for j in range(cols):
            table.cell(i + 1, j).text = str(df_contracts_filtered.iloc[i, j])
    
    ppt_buffer = BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer

# Function to get SQLite connection
def get_sqlite_connection():
    conn = sqlite3.connect('comments.db')
    cursor = conn.cursor()
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS comments (
            id TEXT,
            type TEXT,
            comment TEXT,
            user TEXT,
            timestamp TEXT
        )
    ''')
    conn.commit()
    return conn, cursor

# Main configuration
st.title(t["title"])

# Load data from Supabase
loading_placeholder = st.empty()
with st.spinner(t["loading"]):
    loading_placeholder.write("Chargement de la table purchase_orders...")
    df_po, error_po = load_data("purchase_orders", ["PO_NUMBER", "FOURNISSEUR", "DEPARTEMENT", "MONTANT_EUR", "QUANTITE", "DATE", "TYPE_ACHAT", "STATUT"])
    if error_po:
        loading_placeholder.error(error_po)
        st.stop()
    loading_placeholder.write("Colonnes trouvées dans purchase_orders: " + str(df_po.columns.tolist()))

    loading_placeholder.write("Chargement de la table payment_terms...")
    df_pt, error_pt = load_data("payment_terms", ["FOURNISSEUR", "OLD_DAYS", "NEW_DAYS", "TURNOVER_EUR", "DIVISION", "CONDITION_PAIEMENT", "DELAI_PAIEMENT"])
    if error_pt:
        loading_placeholder.error(error_pt)
        st.stop()
    loading_placeholder.write("Colonnes trouvées dans payment_terms: " + str(df_pt.columns.tolist()))

    loading_placeholder.write("Chargement de la table contracts...")
    df_contracts, error_contracts = load_data("contracts", ["CONTRAT", "FOURNISSEUR", "DATE_EXPIRATION", "MONTANT_MAD", "RESPONSABLE_EMAIL"])
    if error_contracts:
        loading_placeholder.error(error_contracts)
        st.stop()
    loading_placeholder.write("Colonnes trouvées dans contracts: " + str(df_contracts.columns.tolist()))

    # Clear the loading placeholder
    loading_placeholder.empty()

    # Convert data types
    try:
        df_po["DATE"] = pd.to_datetime(df_po["DATE"], errors='coerce')
        df_po["MONTANT_EUR"] = pd.to_numeric(df_po["MONTANT_EUR"], errors='coerce')
        df_po["QUANTITE"] = pd.to_numeric(df_po["QUANTITE"], errors='coerce')
        df_po["STATUT"] = df_po["STATUT"].astype(str)
        df_po["TYPE_ACHAT"] = df_po["TYPE_ACHAT"].astype(str)

        df_pt["NEW_DAYS"] = pd.to_numeric(df_pt["NEW_DAYS"], errors='coerce')
        df_pt["OLD_DAYS"] = pd.to_numeric(df_pt["OLD_DAYS"], errors='coerce')
        df_pt["TURNOVER_EUR"] = pd.to_numeric(df_pt["TURNOVER_EUR"], errors='coerce')
        df_pt["DELAI_PAIEMENT"] = pd.to_numeric(df_pt["DELAI_PAIEMENT"], errors='coerce')

        df_contracts["DATE_EXPIRATION"] = pd.to_datetime(df_contracts["DATE_EXPIRATION"], errors='coerce')
        df_contracts["MONTANT_MAD"] = pd.to_numeric(df_contracts["MONTANT_MAD"], errors='coerce')

        if df_po["DATE"].isna().any() or df_contracts["DATE_EXPIRATION"].isna().any():
            st.warning("Certaines dates n'ont pas pu être converties. Vérifiez le format des données dans Supabase.")
    except Exception as e:
        st.error(f"Erreur lors de la conversion des types de données : {str(e)}")
        st.stop()

st.success(t["data_loaded"])

# Global summary
st.markdown('<div class="section">', unsafe_allow_html=True)
st.subheader(t["summary"])
col_sum1, col_sum2, col_sum3 = st.columns(3)
with col_sum1:
    total_orders = len(df_po)
    st.metric(t["total_orders"], total_orders)
with col_sum2:
    pending_orders = len(df_po[df_po["STATUT"] == "En attente"])
    st.metric(t["pending_orders"], pending_orders)
with col_sum3:
    total_turnover = df_pt["TURNOVER_EUR"].sum()
    st.metric(t["total_turnover"], f"{total_turnover:,.2f} EUR")
st.markdown('</div>', unsafe_allow_html=True)

# Sidebar
with st.sidebar:
    st.header(t["filters_alerts"])
    
    global_search = st.text_input("Recherche globale 🔍", key="global_search")
    
    st.subheader(t["po_filters"])
    fournisseur_po = st.multiselect(t["supplier"] + " (PO)", df_po["FOURNISSEUR"].unique(), default=df_po["FOURNISSEUR"].unique())
    departement = st.multiselect(t["department"], df_po["DEPARTEMENT"].unique(), default=df_po["DEPARTEMENT"].unique())
    type_achat = st.multiselect(t["purchase_type"], df_po["TYPE_ACHAT"].unique(), default=df_po["TYPE_ACHAT"].unique())
    statut = st.multiselect(t["status"], df_po["STATUT"].unique(), default=df_po["STATUT"].unique())
    period = st.slider(t["period"], df_po["DATE"].min().to_pydatetime(), df_po["DATE"].max().to_pydatetime(), 
                       (df_po["DATE"].min().to_pydatetime(), df_po["DATE"].max().to_pydatetime()))

    if not all([fournisseur_po, departement, type_achat, statut]):
        st.warning(t["select_filter"])

    st.subheader(t["pt_filters"])
    fournisseur_pt = st.multiselect(t["supplier"] + " (PT)", df_pt["FOURNISSEUR"].unique(), default=df_pt["FOURNISSEUR"].unique())
    division = st.multiselect(t["division"], df_pt["DIVISION"].unique(), default=df_pt["DIVISION"].unique())
    period_pt = st.slider(t["period"], datetime(2023, 1, 1), datetime(2025, 12, 31), (datetime(2023, 1, 1), datetime(2025, 12, 31)))

    if not all([fournisseur_pt, division]):
        st.warning(t["select_filter"])

    st.subheader(t["contract_filters"])
    fournisseur_contract = st.multiselect(t["supplier"] + " (Contrats)", df_contracts["FOURNISSEUR"].unique(), default=df_contracts["FOURNISSEUR"].unique())
    expiration_period = st.slider(t["period"], df_contracts["DATE_EXPIRATION"].min().to_pydatetime(), 
                                  df_contracts["DATE_EXPIRATION"].max().to_pydatetime(), 
                                  (df_contracts["DATE_EXPIRATION"].min().to_pydatetime(), df_contracts["DATE_EXPIRATION"].max().to_pydatetime()))

    st.subheader(t["alerts"])
    seuil_alert = st.number_input(t["amount_threshold"], min_value=0.0, value=100000.0, step=1000.0)
    seuil_delai = st.number_input(t["delay_threshold"], min_value=0, value=5, step=1)

    df_po_filtered = df_po[
        (df_po["FOURNISSEUR"].isin(fournisseur_po)) &
        (df_po["DEPARTEMENT"].isin(departement)) &
        (df_po["TYPE_ACHAT"].isin(type_achat)) &
        (df_po["STATUT"].isin(statut)) &
        (df_po["DATE"].between(period[0], period[1]))
    ]
    df_pt_filtered = df_pt[
        (df_pt["FOURNISSEUR"].isin(fournisseur_pt)) &
        (df_pt["DIVISION"].isin(division))
    ]
    df_contracts_filtered = df_contracts[
        (df_contracts["FOURNISSEUR"].isin(fournisseur_contract)) &
        (df_contracts["DATE_EXPIRATION"].between(expiration_period[0], expiration_period[1]))
    ]

    if global_search:
        df_po_filtered = df_po_filtered[df_po_filtered.apply(lambda row: global_search.lower() in str(row).lower(), axis=1)]
        df_pt_filtered = df_pt_filtered[df_pt_filtered.apply(lambda row: global_search.lower() in str(row).lower(), axis=1)]
        df_contracts_filtered = df_contracts_filtered[df_contracts_filtered.apply(lambda row: global_search.lower() in str(row).lower(), axis=1)]

    alerts_displayed = False
    for index, row in df_po_filtered.iterrows():
        if row["MONTANT_EUR"] > seuil_alert:
            st.warning(f"⚠️ {row['FOURNISSEUR']}: {row['MONTANT_EUR']:,.2f} EUR exceeds threshold", icon="🚨")
            alerts_displayed = True
        if row["STATUT"] == "En attente":
            st.warning(f"⚠️ Order {row['PO_NUMBER']} pending", icon="⏳")
            alerts_displayed = True
    for index, row in df_contracts_filtered.iterrows():
        days_left = (row["DATE_EXPIRATION"] - pd.Timestamp.now()).days
        if days_left <= 60:
            st.warning(f"⚠️ Contract {row['CONTRAT']} expires in {days_left} days", icon="📅")
            alerts_displayed = True
    for index, row in df_pt_filtered.iterrows():
        if row["DELAI_PAIEMENT"] > seuil_delai:
            st.warning(f"⚠️ {row['FOURNISSEUR']}: {row['DELAI_PAIEMENT']} days payment delay", icon="⏰")
            alerts_displayed = True
    if not alerts_displayed:
        st.info(t["no_alerts"])

    if smtp_available and st.button(t["check_alerts"], key="check_alerts_btn"):
        notifications_sent = check_alerts(df_contracts_filtered, df_po_filtered, df_pt_filtered)
        if notifications_sent > 0:
            st.success(f"{notifications_sent} {t['alerts'].lower()}(s) sent.")
        else:
            st.info("Aucune alerte à notifier.")

    st.subheader(t["help"])
    st.write(t["help_text"])

# Main tabs
tab1, tab2, tab3 = st.tabs([t["po_tab"], t["pt_tab"], t["contract_tab"]])

with tab1:
    st.markdown('<div class="section">', unsafe_allow_html=True)
    st.subheader(t["po_header"])
    if df_po_filtered.empty:
        st.warning(t["no_data"])
    else:
        fig_po_count = None
        fig_status = None
        fig_type = None
        fig_monthly = None
        fig_annual = None

        col1, col2 = st.columns(2)
        with col1:
            st.subheader(t["po_by_dept"])
            view = st.radio("View", ["Monthly", "Annual"], key="po_view")
            if view == "Monthly":
                df_grouped = df_po_filtered.groupby([df_po_filtered["DATE"].dt.to_period("M").astype(str), "DEPARTEMENT"]).agg({"MONTANT_EUR": "sum"}).reset_index()
                fig_po_count = px.bar(df_grouped, x="DATE", y="MONTANT_EUR", color="DEPARTEMENT", title=t["po_by_dept"], text="MONTANT_EUR", color_discrete_sequence=color_schemes[color_scheme])
                fig_po_count.update_traces(textposition="outside")
                fig_po_count.update_layout(xaxis_title="Month", yaxis_title="Amount (EUR)")
            else:
                df_grouped = df_po_filtered.groupby("DEPARTEMENT").agg({"MONTANT_EUR": "sum"}).reset_index()
                fig_po_count = px.bar(df_grouped, x="DEPARTEMENT", y="MONTANT_EUR", color="DEPARTEMENT", title=t["po_by_dept"], text="MONTANT_EUR", color_discrete_sequence=color_schemes[color_scheme])
                fig_po_count.update_traces(textposition="outside")
                fig_po_count.update_layout(showlegend=False)
            st.plotly_chart(fig_po_count, use_container_width=True)

        with col2:
            st.subheader(t["amount_quantity"])
            view = st.radio("View", ["Monthly", "Annual"], key="amount_view")
            if view == "Monthly":
                df_monthly = df_po_filtered.groupby(df_po_filtered["DATE"].dt.to_period("M").astype(str)).agg({"MONTANT_EUR": "sum", "QUANTITE": "sum"}).reset_index()
                fig_monthly = px.bar(df_monthly, x="DATE", y="MONTANT_EUR", title=t["amount_quantity"], text="MONTANT_EUR", color_discrete_sequence=color_schemes[color_scheme])
                fig_monthly.update_traces(textposition="outside")
                fig_monthly.update_layout(height=400)
            else:
                df_annual = df_po_filtered.groupby(df_po_filtered["DATE"].dt.year).agg({"MONTANT_EUR": "sum", "QUANTITE": "sum"}).reset_index()
                fig_annual = px.bar(df_annual, x="DATE", y="MONTANT_EUR", title=t["amount_quantity"], text="MONTANT_EUR", color_discrete_sequence=color_schemes[color_scheme])
                fig_annual.update_traces(textposition="outside")
            fig_to_plot = fig_monthly if view == "Monthly" else fig_annual
            st.plotly_chart(fig_to_plot, use_container_width=True)

        st.subheader(t["status_dist"])
        status_counts = df_po_filtered.groupby("STATUT").size().reset_index(name="Count")
        fig_status = px.pie(status_counts, names="STATUT", values="Count", title=t["status_dist"], hole=0.4, color_discrete_sequence=color_schemes[color_scheme])
        st.plotly_chart(fig_status, use_container_width=True)

        st.subheader(t["type_dist"])
        fig_type = px.pie(df_po_filtered.groupby("TYPE_ACHAT").size().reset_index(name="Count"), names="TYPE_ACHAT", values="Count", title=t["type_dist"], hole=0.4, color_discrete_sequence=color_schemes[color_scheme])
        st.plotly_chart(fig_type, use_container_width=True)

        st.subheader(t["forecast"])
        predict_by = st.selectbox(t["predict_by"], ["Département", "Fournisseur"], key="predict_by")
        if predict_by == "Département":
            options = df_po_filtered["DEPARTEMENT"].unique()
        else:
            options = df_po_filtered["FOURNISSEUR"].unique()
        selected_option = st.selectbox(f"Sélectionner {predict_by.lower()}", options, key="predict_option")

        column_name = "DEPARTEMENT" if predict_by == "Département" else "FOURNISSEUR"
        df_predict = df_po_filtered[df_po_filtered[column_name] == selected_option]

        if df_predict.empty:
            st.warning(f"Aucune donnée disponible pour {predict_by.lower()} '{selected_option}'. Vérifiez les filtres ou les données dans Supabase.")
        else:
            df_predict = df_predict.groupby(df_predict["DATE"].dt.to_period("M")).agg({"MONTANT_EUR": "sum"}).reset_index()
            df_predict["DATE"] = df_predict["DATE"].dt.to_timestamp()

            if df_predict["DATE"].isna().any():
                st.error(f"Données de date invalides pour {predict_by.lower()} '{selected_option}'. Vérifiez le format des dates dans la table 'purchase_orders'.")
            else:
                df_predict["time_index"] = (df_predict["DATE"] - df_predict["DATE"].min()).dt.days

                st.write(f"Données agrégées pour {predict_by.lower()} '{selected_option}' :")
                st.dataframe(df_predict[["DATE", "MONTANT_EUR"]])

                if len(df_predict) < 2:
                    st.info(f"Données insuffisantes pour une prévision (un seul mois disponible pour {predict_by.lower()} '{selected_option}'). Affichage des données existantes.")
                    fig_predict = px.line(df_predict, x="DATE", y="MONTANT_EUR", title=f"{t['forecast']} pour {selected_option}", color_discrete_sequence=color_schemes[color_scheme])
                    st.plotly_chart(fig_predict, use_container_width=True)
                else:
                    X = df_predict[["time_index"]]  # Keep as DataFrame
                    y = df_predict["MONTANT_EUR"]
                    model = LinearRegression()
                    model.fit(X, y)
                    future_dates = pd.date_range(start=df_predict["DATE"].max() + pd.offsets.MonthBegin(1), periods=6, freq="MS")
                    future_time_index = pd.DataFrame([(date - df_predict["DATE"].min()).days for date in future_dates], columns=["time_index"])
                    future_predictions = model.predict(future_time_index)
                    fig_predict = px.line(df_predict, x="DATE", y="MONTANT_EUR", title=f"{t['forecast']} pour {selected_option}", color_discrete_sequence=color_schemes[color_scheme])
                    fig_predict.add_scatter(x=future_dates, y=future_predictions, mode="lines+markers", name="Prévision", line=dict(dash="dash"))
                    st.plotly_chart(fig_predict, use_container_width=True)

        st.subheader(t["compare_fournisseurs"])
        fournisseurs_compare = st.multiselect(t["supplier"], df_po_filtered["FOURNISSEUR"].unique(), default=df_po_filtered["FOURNISSEUR"].unique()[:3], key="compare_fournisseurs")
        if fournisseurs_compare:
            df_compare = df_po_filtered[df_po_filtered["FOURNISSEUR"].isin(fournisseurs_compare)]
            df_compare = df_compare.groupby("FOURNISSEUR").agg({
                "MONTANT_EUR": "sum",
                "QUANTITE": "sum",
                "PO_NUMBER": "count"
            }).reset_index()
            df_compare["Taux_Pending"] = df_compare["FOURNISSEUR"].apply(
                lambda x: len(df_po_filtered[(df_po_filtered["FOURNISSEUR"] == x) & (df_po_filtered["STATUT"] == "En attente")]) / 
                          len(df_po_filtered[df_po_filtered["FOURNISSEUR"] == x]) * 100 if len(df_po_filtered[df_po_filtered["FOURNISSEUR"] == x]) > 0 else 0
            )
            df_compare = df_compare.merge(df_pt_filtered[["FOURNISSEUR", "NEW_DAYS"]], on="FOURNISSEUR", how="left")
            for col in ["MONTANT_EUR", "QUANTITE", "Taux_Pending", "NEW_DAYS"]:
                df_compare[col] = (df_compare[col] - df_compare[col].min()) / (df_compare[col].max() - df_compare[col].min() + 1e-6)
            
            fig_radar = go.Figure()
            for fournisseur in fournisseurs_compare:
                df_fournisseur = df_compare[df_compare["FOURNISSEUR"] == fournisseur]
                fig_radar.add_trace(go.Scatterpolar(
                    r=df_fournisseur[["MONTANT_EUR", "QUANTITE", "Taux_Pending", "NEW_DAYS"]].values.flatten().tolist() + [df_fournisseur["MONTANT_EUR"].iloc[0]],
                    theta=["Montant", "Quantité", "Taux Pending", "Délai Paiement", "Montant"],
                    fill="toself",
                    name=fournisseur
                ))
            fig_radar.update_layout(polar=dict(radialaxis=dict(visible=True, range=[0, 1])), showlegend=True, title=t["compare_fournisseurs"])
            st.plotly_chart(fig_radar, use_container_width=True)

        st.subheader(t["reorder"])
        type_achat_reorder = st.multiselect(t["purchase_type"], df_po_filtered["TYPE_ACHAT"].unique(), key="reorder_type")
        threshold = st.number_input(t["reorder_threshold"], min_value=0, value=100, step=10)
        if type_achat_reorder:
            df_reorder = df_po_filtered[df_po_filtered["TYPE_ACHAT"].isin(type_achat_reorder)].groupby("TYPE_ACHAT").agg({"QUANTITE": "sum"}).reset_index()
            df_reorder["Suggestion"] = df_reorder["QUANTITE"].apply(lambda x: t["reorder"] if x < threshold else "Stock suffisant")
            st.dataframe(df_reorder)

        st.subheader("Purchase Orders Details")
        search_term = st.text_input("Search PO", "", key="po_search")
        filtered_df = df_po_filtered[["PO_NUMBER", "FOURNISSEUR", "DEPARTEMENT", "MONTANT_EUR", "QUANTITE", "DATE", "STATUT"]]
        if search_term:
            filtered_df = filtered_df[filtered_df.apply(lambda row: search_term.lower() in str(row).lower(), axis=1)]
        gb = GridOptionsBuilder.from_dataframe(filtered_df)
        gb.configure_pagination(paginationAutoPageSize=True)
        gb.configure_side_bar()
        gb.configure_default_column(editable=True, groupable=True)
        grid_options = gb.build()
        AgGrid(filtered_df, gridOptions=grid_options, height=200, width='100%', fit_columns_on_grid_load=True)

        st.subheader(t["comments"])
        selected_po = st.text_input("PO_NUMBER", key="comment_po_number")
        comment = st.text_area(t["comment_text"], key="comment_text")
        user = st.text_input(t["comment_user"], key="comment_user")
        if st.button(t["add_comment"], key="add_comment_po"):
            if comment and user and selected_po:
                conn, cursor = get_sqlite_connection()
                try:
                    cursor.execute("INSERT INTO comments (id, type, comment, user, timestamp) VALUES (?, ?, ?, ?, ?)", 
                                   (selected_po, "PO", comment, user, datetime.now().strftime('%Y-%m-%d %H:%M:%S')))
                    conn.commit()
                    st.success("Commentaire ajouté !")
                finally:
                    conn.close()
        conn, cursor = get_sqlite_connection()
        try:
            comments_df = pd.read_sql_query("SELECT * FROM comments WHERE id = ? AND type = ?", conn, params=(selected_po, "PO"))
            st.dataframe(comments_df[["comment", "user", "timestamp"]], use_container_width=True)
        finally:
            conn.close()

    st.markdown('</div>', unsafe_allow_html=True)

with tab2:
    st.markdown('<div class="section">', unsafe_allow_html=True)
    st.subheader(t["pt_header"])
    if df_pt_filtered.empty:
        st.warning(t["no_data"])
    else:
        fig_new_terms = None
        fig_old_terms = None

        col1, col2 = st.columns(2)
        with col1:
            st.subheader(t["new_terms"])
            fig_new_terms = px.pie(
                df_pt_filtered.groupby(pd.cut(df_pt_filtered["NEW_DAYS"], bins=[0, 45, 60, float("inf")],
                                              labels=["≤45 days", "45-60 days", "≥60 days"])).size().reset_index(name="Count"),
                names="NEW_DAYS",
                values="Count",
                hole=0.4,
                title=t["new_terms"],
                color_discrete_sequence=color_schemes[color_scheme]
            )
            st.plotly_chart(fig_new_terms, use_container_width=True)

        with col2:
            st.subheader(t["old_terms"])
            fig_old_terms = px.pie(
                df_pt_filtered.groupby(pd.cut(df_pt_filtered["OLD_DAYS"], bins=[0, 45, 60, float("inf")],
                                              labels=["≤45 days", "45-60 days", "≥60 days"])).size().reset_index(name="Count"),
                names="OLD_DAYS",
                values="Count",
                hole=0.4,
                title=t["old_terms"],
                color_discrete_sequence=color_schemes[color_scheme]
            )
            st.plotly_chart(fig_old_terms, use_container_width=True)

        st.subheader(t["heatmap"])
        metric = st.selectbox("Métrique", ["Turnover (EUR)", "Délai Paiement (jours)"], key="heatmap_metric")
        df_heatmap = df_pt_filtered.groupby(["FOURNISSEUR", "DIVISION"]).agg({
            "TURNOVER_EUR": "sum",
            "NEW_DAYS": "mean"
        }).reset_index()
        if metric == "Turnover (EUR)":
            z = df_heatmap.pivot(index="FOURNISSEUR", columns="DIVISION", values="TURNOVER_EUR").fillna(0)
        else:
            z = df_heatmap.pivot(index="FOURNISSEUR", columns="DIVISION", values="NEW_DAYS").fillna(0)
        fig_heatmap = px.imshow(z, title=f"{t['heatmap']} ({metric})", color_continuous_scale=color_schemes[color_scheme])
        st.plotly_chart(fig_heatmap, use_container_width=True)

        st.subheader(t["kpis"])
        col_kpi1, col_kpi2, col_kpi3 = st.columns(3)
        with col_kpi1:
            turnover_total = df_pt_filtered["TURNOVER_EUR"].sum()
            st.metric(t["turnover"], f"{turnover_total:,.2f}")
        with col_kpi2:
            improved_suppliers = len(df_pt_filtered[df_pt_filtered["NEW_DAYS"] < df_pt_filtered["OLD_DAYS"]])
            total_suppliers = len(df_pt_filtered)
            pt_improvement = (improved_suppliers / total_suppliers) * 100 if total_suppliers > 0 else 0
            st.metric(t["improvement"], f"{pt_improvement:.2f}")
        with col_kpi3:
            cash_flow_gain = ((df_pt_filtered["OLD_DAYS"] - df_pt_filtered["NEW_DAYS"]) * df_pt_filtered["TURNOVER_EUR"] / 360).sum()
            st.metric(t["cash_flow"], f"{cash_flow_gain:,.2f}")

        st.subheader(t["terms_by_division"])
        df_division = df_pt_filtered.groupby("DIVISION").agg({"NEW_DAYS": "mean", "OLD_DAYS": "mean", "TURNOVER_EUR": "sum"}).reset_index()
        fig_division = px.bar(df_division, x="DIVISION", y=["NEW_DAYS", "OLD_DAYS"], barmode="group", title=t["terms_by_division"], color_discrete_sequence=color_schemes[color_scheme])
        fig_division.add_scatter(x=df_division["DIVISION"], y=df_division["TURNOVER_EUR"], mode="lines+markers", name="Turnover", yaxis="y2")
        fig_division.update_layout(yaxis2=dict(title="Turnover (EUR)", overlaying="y", side="right"), yaxis_title="Days")
        st.plotly_chart(fig_division, use_container_width=True)

        st.subheader("KPI by Division")
        kpi_df = df_pt_filtered.groupby("DIVISION").agg({"TURNOVER_EUR": "sum", "NEW_DAYS": "mean", "OLD_DAYS": "mean"}).reset_index()
        kpi_df["Improvement"] = ((kpi_df["OLD_DAYS"] - kpi_df["NEW_DAYS"]) / kpi_df["OLD_DAYS"] * 100).round(2)
        st.dataframe(kpi_df, use_container_width=True)

    st.markdown('</div>', unsafe_allow_html=True)

with tab3:
    st.markdown('<div class="section">', unsafe_allow_html=True)
    st.subheader(t["contract_header"])
    if df_contracts_filtered.empty:
        st.warning(t["no_data"])
    else:
        gb = GridOptionsBuilder.from_dataframe(df_contracts_filtered)
        gb.configure_pagination(paginationAutoPageSize=True)
        gb.configure_side_bar()
        gb.configure_default_column(editable=False, groupable=True)
        grid_options = gb.build()
        AgGrid(df_contracts_filtered, gridOptions=grid_options, height=200, width='100%', fit_columns_on_grid_load=True)

        if smtp_available and st.button(t["send_contract_reminders"], key="send_contract_reminders"):
            notifications_sent = 0
            for _, row in df_contracts_filtered.iterrows():
                days_left = (row["DATE_EXPIRATION"] - pd.Timestamp.now()).days
                if days_left <= 60 and pd.notna(row["RESPONSABLE_EMAIL"]):
                    subject = f"Rappel : Contrat {row['CONTRAT']} expire bientôt"
                    body = f"Le contrat {row['CONTRAT']} avec {row['FOURNISSEUR']} expire dans {days_left} jours.\nDate d'expiration : {row['DATE_EXPIRATION'].strftime('%d/%m/%Y')}"
                    if send_email(subject, body, row["RESPONSABLE_EMAIL"]):
                        notifications_sent += 1
            if notifications_sent > 0:
                st.success(f"{notifications_sent} rappel(s) envoyé(s).")
        else:
            st.info("Aucun rappel à envoyer.")

        st.subheader(t["comments"])
        selected_contract = st.text_input("CONTRAT", key="comment_contract_number")
        comment = st.text_area(t["comment_text"], key="comment_text_contract")
        user = st.text_input(t["comment_user"], key="comment_user_contract")
        if st.button(t["add_comment"], key="add_comment_contract"):
            if comment and user and selected_contract:
                conn, cursor = get_sqlite_connection()
                try:
                    cursor.execute("INSERT INTO comments (id, type, comment, user, timestamp) VALUES (?, ?, ?, ?, ?)", 
                                   (selected_contract, "Contract", comment, user, datetime.now().strftime('%Y-%m-%d %H:%M:%S')))
                    conn.commit()
                    st.success("Commentaire ajouté !")
                finally:
                    conn.close()
        conn, cursor = get_sqlite_connection()
        try:
            comments_df = pd.read_sql_query("SELECT * FROM comments WHERE id = ? AND type = ?", conn, params=(selected_contract, "Contract"))
            st.dataframe(comments_df[["comment", "user", "timestamp"]], use_container_width=True)
        finally:
            conn.close()

    st.markdown('</div>', unsafe_allow_html=True)

# Export PowerPoint
st.markdown('<div class="section">', unsafe_allow_html=True)
if st.button(t["export_ppt"], key="export_ppt_btn"):
    figs = {
        "po_by_dept": fig_po_count,
        "status_dist": fig_status,
        "type_dist": fig_type,
        "new_terms": fig_new_terms,
        "old_terms": fig_old_terms
    }
    ppt_buffer = export_to_ppt(df_po_filtered, df_pt_filtered, df_contracts_filtered, figs)
    st.download_button(label=t["export_ppt"], data=ppt_buffer.getvalue(), file_name="dashboard_report.pptx", 
                       mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", key="download_ppt")
st.markdown('</div>', unsafe_allow_html=True)

# Footer
st.markdown("---")
st.markdown(f"{t['footer']} {datetime.now().strftime('%d/%m/%Y %H:%M')}")
